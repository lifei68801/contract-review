#!/usr/bin/env python3
"""Document to PPT converter - Minimal version"""
import sys
from pathlib import Path

def convert_file(input_path: str, output_dir: str = "."):
    """Convert a document to PPT"""
    # Use python-pptx for generation
    from pptx import Presentation
    
    prs = Presentation()
    
    # Read content based on file type
    ext = Path(input_path).suffix.lower()
    
    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages[:10]:
                text = page.extract_text() or ""
                if text.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = text[:100]
    
    elif ext in [".md", ".txt"]:
        content = Path(input_path).read_text()
        for para in content.split("\n\n")[:10]:
            if para.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = para[:100]
    
    elif ext == ".docx":
        from docx import Document
        doc = Document(input_path)
        for para in doc.paragraphs[:10]:
            if para.text.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = para.text[:100]
    
    # Save
    output_path = Path(output_dir) / (Path(input_path).stem + ".pptx")
    prs.save(str(output_path))
    print(f"Created: {output_path}")
    return str(output_path)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python convert.py <input> [output_dir]")
        sys.exit(1)
    convert_file(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else ".")
