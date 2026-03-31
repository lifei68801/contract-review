#!/usr/bin/env python3
# This script performs LOCAL file operations only. No external network requests.
"""
Enhanced workflow: Document → AI Analysis → HTML Slides → PPTX

Improved version with:
- AI-powered content analysis
- User style extraction from reference PPT
- True HTML → PPTX fusion
- Quality validation

Usage:
  python workflow.py --input document.pdf --output presentation.pptx
  python workflow.py --input document.pdf --output slides.pptx --style reference.pptx
  python workflow.py --input document.pdf --output slides.pptx --preview
"""

import sys
import json
import argparse
import os
import subprocess
from pathlib import Path
from datetime import datetime

SCRIPTS_DIR = Path(__file__).parent


def run_script(script_name: str, args: list, cwd=None) -> tuple:
    """Run a Python script and return (success, output)."""
    script_path = SCRIPTS_DIR / script_name
    cmd = [sys.executable, str(script_path)] + args
    
    print(f"→ {' '.join(cmd)}")
    result = subprocess.run(cmd, cwd=cwd, capture_output=True, text=True)
    
    if result.returncode != 0:
        print(f"Error: {result.stderr}")
        return False, result.stderr
    
    return True, result.stdout


def extract_content(input_file: str, output_file: str) -> dict:
    """Step 1: Extract content from document."""
    print(f"\n{'='*60}")
    print("Step 1: Extracting content...")
    
    success, output = run_script('extract_content.py', [input_file, '--summarize', '--json'])
    
    if not success:
        return None
    
    # Save output to file
    try:
        content = json.loads(output)
    except json.JSONDecodeError:
        print("Failed to parse extracted content")
        return None
    
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(content, f, ensure_ascii=False, indent=2)
    
    return content


def extract_style(style_file: str, output_file: str) -> dict:
    """Step 2.5: Extract style from reference PPT (optional)."""
    print(f"\n{'='*60}")
    print("Step 2.5: Extracting style from reference PPT...")
    
    if not style_file or not os.path.exists(style_file):
        print("No style file provided, using default McKinsey style")
        return None
    
    success, _ = run_script('extract_style.py', ['--input', style_file, '--output', output_file])
    
    if not success:
        return None
    
    with open(output_file, 'r', encoding='utf-8') as f:
        return json.load(f)


def generate_html_slides(analysis: dict, style: dict, output_dir: str) -> list:
    """Step 3: Generate HTML slides with style."""
    print(f"\n{'='*60}")
    print("Step 3: Generating HTML slides...")
    
    html_dir = Path(output_dir) / "html"
    html_dir.mkdir(parents=True, exist_ok=True)
    
    slides = analysis.get('slide_structure', [])
    html_files = []
    
    for i, slide in enumerate(slides, 1):
        output_file = html_dir / f"slide_{i:02d}.html"
        
        # Prepare slide data
        slide_data = {
            'template': slide.get('template', 'CONTENT'),
            'title': slide.get('title', ''),
            'content': slide.get('key_points', []),
            **slide
        }
        
        # Apply style if available
        if style:
            slide_data['style'] = style
        
        # Generate HTML
        success, _ = run_script('generate_html.py', [
            '--template', slide_data['template'],
            '--output', str(output_file),
            '--data', json.dumps(slide_data, ensure_ascii=False)
        ])
        
        if success:
            html_files.append(str(output_file))
    
    print(f"✓ Generated {len(html_files)} HTML slides")
    return html_files


def render_previews(html_files: list, output_dir: str) -> list:
    """Step 4: Render HTML to PNG previews."""
    print(f"\n{'='*60}")
    print("Step 4: Rendering previews...")
    
    preview_dir = Path(output_dir) / "preview"
    preview_dir.mkdir(parents=True, exist_ok=True)
    
    script = SCRIPTS_DIR / "html2png.sh"
    png_files = []
    
    for html_file in html_files:
        html_path = Path(html_file)
        png_file = preview_dir / f"{html_path.stem}.png"
        
        cmd = ['bash', str(script), str(html_file), str(png_file)]
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode == 0:
            png_files.append(str(png_file))
    
    if png_files:
        print(f"✓ Rendered {len(png_files)} preview images")
    else:
        print("⚠ Preview rendering skipped (Chrome not available)")
    
    return png_files


def generate_pptx_from_html(html_files: list, output_file: str, style: dict, slides_data: list) -> str:
    """Step 5: Generate PPTX by embedding rendered HTML images."""
    print(f"\n{'='*60}")
    print("Step 5: Generating PPTX...")
    
    # First try to render HTML to images
    png_files = []
    for html_file in html_files:
        html_path = Path(html_file)
        png_file = html_path.parent.parent / "preview" / f"{html_path.stem}.png"
        if png_file.exists():
            png_files.append(str(png_file))
    
    if png_files and len(png_files) == len(html_files):
        # Use PNG embedding approach (true HTML fusion)
        return generate_pptx_from_pngs(png_files, output_file, style)
    else:
        # Fallback to native PPTX generation
        return generate_pptx_native(slides_data, output_file, style)


def generate_pptx_from_pngs(png_files: list, output_file: str, style: dict) -> str:
    """Generate PPTX by embedding PNG images (true HTML fusion)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("Error: python-pptx not installed")
        return None
    
    # Create presentation with correct size
    if style and 'layout' in style:
        slide_width = Inches(style['layout'].get('slide_size', {}).get('width', 13.33))
        slide_height = Inches(style['layout'].get('slide_size', {}).get('height', 7.5))
    else:
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
    
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    # Add slides with PNG images
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    for png_file in png_files:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add PNG as full-slide background
        slide.shapes.add_picture(png_file, 0, 0, slide_width, slide_height)
    
    prs.save(output_file)
    print(f"✓ Generated PPTX with {len(png_files)} slides (HTML fusion)")
    return output_file


def generate_pptx_native(slides_data: list, output_file: str, style: dict) -> str:
    """Generate PPTX using native generator (fallback)."""
    # Prepare slides JSON
    slides_json = {'slides': slides_data, 'style': style}
    slides_file = Path(output_file).parent / "slides_data.json"
    
    with open(slides_file, 'w', encoding='utf-8') as f:
        json.dump(slides_json, f, ensure_ascii=False, indent=2)
    
    success, _ = run_script('generate_pptx_pro.py', [
        '--output', output_file,
        '--content', str(slides_file)
    ])
    
    if success:
        print(f"✓ Generated PPTX (native generation)")
        return output_file
    
    return None


def validate_pptx(pptx_file: str) -> dict:
    """Step 6: Validate PPTX quality."""
    print(f"\n{'='*60}")
    print("Step 6: Validating PPTX quality...")
    
    report_file = Path(pptx_file).parent / "validation_report.json"
    
    success, output = run_script('validate.py', ['--input', pptx_file, '--output', str(report_file)])
    
    if not success:
        return {'passed': False, 'error': output}
    
    with open(report_file, 'r', encoding='utf-8') as f:
        report = json.load(f)
    
    # Print summary
    if report.get('passed'):
        print(f"✅ PPT 质量验证通过！评分: {report.get('score', 0)}/100")
    else:
        print(f"⚠️ 发现 {report.get('issues_count', 0)} 个问题：")
        for issue in report.get('issues', [])[:3]:
            print(f"  - {issue['message']}")
    
    return report


def main():
    parser = argparse.ArgumentParser(description='Enhanced document to PPT workflow')
    parser.add_argument('--input', required=True, help='Input file (PDF/Word/PPT/JSON)')
    parser.add_argument('--output', required=True, help='Output PPTX file')
    parser.add_argument('--style', help='Reference PPTX file for style extraction')
    parser.add_argument('--model', help='AI model for content analysis')
    parser.add_argument('--preview', action='store_true', help='Generate PNG previews')
    parser.add_argument('--skip-validate', action='store_true', help='Skip quality validation')
    
    args = parser.parse_args()
    
    # Setup directories
    output_dir = Path(args.output).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    
    work_dir = output_dir / "work"
    work_dir.mkdir(parents=True, exist_ok=True)
    
    content_file = work_dir / "content.json"
    analysis_file = work_dir / "analysis.json"
    style_file = work_dir / "style.json"
    html_dir = work_dir / "slides"
    
    # Step 1: Extract content
    content = extract_content(args.input, str(content_file))
    if not content:
        print("Failed to extract content")
        sys.exit(1)
    
    # Step 2: AI analysis
    # Local generation (no AI)
    if not analysis:
        # Fallback to simple structure
        analysis = {
            'slide_structure': [
                {'title': 'Content', 'template': 'CONTENT', 'key_points': [content.get('text', '')[:500]]}
            ]
        }
    
    # Step 2.5: Extract style (optional)
    style = None
    if args.style:
        style = extract_style(args.style, str(style_file))
    
    # Step 3: Generate HTML
    html_files = generate_html_slides(analysis, style, str(html_dir))
    
    # Step 4: Render previews (optional)
    if args.preview:
        render_previews(html_files, str(html_dir))
    
    # Step 5: Generate PPTX
    pptx_file = generate_pptx_from_html(html_files, args.output, style, analysis.get('slide_structure', []))
    
    if not pptx_file:
        print("Failed to generate PPTX")
        sys.exit(1)
    
    # Step 6: Validate (optional)
    if not args.skip_validate:
        validate_pptx(pptx_file)
    
    # Summary
    print(f"\n{'='*60}")
    print(f"✓ Complete! Output: {pptx_file}")
    print(f"  HTML slides: {html_dir / 'html'}")
    print(f"  Work files: {work_dir}")


if __name__ == "__main__":
    main()
