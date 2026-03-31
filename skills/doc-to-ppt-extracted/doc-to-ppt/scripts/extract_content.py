#!/usr/bin/env python3
"""
Extract text content from PDF, Word, and PowerPoint files.
Usage: python extract_content.py <file_path>
Output: JSON with extracted text and metadata
"""

import sys
import json
import os

def extract_pdf(file_path):
    """Extract text from PDF using pdfplumber"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(f"--- Page {i+1} ---\n{text}")
        return {"success": True, "text": "\n\n".join(text_parts), "pages": len(pdf.pages)}
    except ImportError:
        # Fallback to PyPDF2
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        text_parts.append(f"--- Page {i+1} ---\n{text}")
            return {"success": True, "text": "\n\n".join(text_parts), "pages": len(reader.pages)}
        except ImportError:
            return {"success": False, "error": "No PDF library installed. Run: pip install pdfplumber or PyPDF2"}
    except Exception as e:
        return {"success": False, "error": str(e)}

def extract_word(file_path):
    """Extract text from Word document using python-docx"""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                text_parts.append("\n[Table]\n" + "\n".join(table_text))
        return {"success": True, "text": "\n\n".join(text_parts), "paragraphs": len(doc.paragraphs)}
    except ImportError:
        return {"success": False, "error": "python-docx not installed. Run: pip install python-docx"}
    except Exception as e:
        return {"success": False, "error": str(e)}

def extract_ppt(file_path):
    """Extract text from PowerPoint using python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                slides_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
        return {
            "success": True, 
            "text": "\n\n".join(slides_content), 
            "slides": len(prs.slides),
            "slide_layouts": [slide.slide_layout.name for slide in prs.slides]
        }
    except ImportError:
        return {"success": False, "error": "python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e:
        return {"success": False, "error": str(e)}

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Usage: python extract_content.py <file_path>"}))
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"success": False, "error": f"File not found: {file_path}"}))
        sys.exit(1)
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        result = extract_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        result = extract_word(file_path)
    elif ext in ['.pptx', '.ppt']:
        result = extract_ppt(file_path)
    else:
        result = {"success": False, "error": f"Unsupported file type: {ext}"}
    
    result["file"] = file_path
    result["type"] = ext
    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
